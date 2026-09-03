#!/usr/bin/env python3
"""
Generates the Executive PowerPoint Presentation for the School of Integrated Engineering:
'Academic Resourcing & Faculty Workload Architecture'
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Paths
ARTIFACT_DIR = r"C:\Users\Michael.Richards\.gemini\antigravity\brain\f996c3e8-5451-419a-bf3c-bda58f5fe43f"
IMG_CARRYING = os.path.join(ARTIFACT_DIR, "boulder_joint_carrying_1788367912553.jpg")
IMG_FACETS = os.path.join(ARTIFACT_DIR, "boulder_divided_facets_1788367928243.jpg")
IMG_PIECES = os.path.join(ARTIFACT_DIR, "boulder_broken_pieces_1788367943230.jpg")
IMG_MULTITASK = os.path.join(ARTIFACT_DIR, "faculty_multitask_roles_1788367958751.jpg")
IMG_SPOTLIGHT = os.path.join(ARTIFACT_DIR, "teaching_load_spotlight_1788367976549.jpg")

OUTPUT_PPTX = "School_of_Engineering_Workload_Architecture.pptx"

# Color Palette
NAVY_DARK = RGBColor(15, 30, 54)        # #0F1E36 (Headers / Dark Cards)
NAVY_MID = RGBColor(30, 58, 138)        # #1E3A8A (Brand Primary)
BLUE_ACCENT = RGBColor(2, 132, 199)     # #0284C7 (Subtitles / Highlights)
BLUE_LIGHT = RGBColor(224, 242, 254)    # #E0F2FE (Badge Backgrounds)
AMBER_ACCENT = RGBColor(217, 119, 6)    # #D97706 (Alerts / Highlights)
AMBER_LIGHT = RGBColor(254, 243, 199)   # #FEF3C7 (Alert BGs)
GREEN_ACCENT = RGBColor(16, 185, 129)   # #10B981 (Success / Positive)
BG_LIGHT = RGBColor(248, 250, 252)      # #F8FAFC (Slide BG)
CARD_BG = RGBColor(255, 255, 255)       # Pure White
BORDER_COLOR = RGBColor(226, 232, 240)  # Light Slate Border
TEXT_MAIN = RGBColor(30, 41, 59)        # #1E293B Body text
TEXT_MUTED = RGBColor(100, 116, 139)    # #64748B Subdued text


def apply_header(slide, title_text, subtitle_text, category="SCHOOL OF INTEGRATED ENGINEERING"):
    """Adds a standardized, clean executive header to a content slide."""
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = BLUE_ACCENT

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.55))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY_DARK

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED


def set_speaker_notes(slide, notes_text):
    """Sets speaker notes for the slide."""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text.strip()


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Title & Vision - Carrying the Load Together
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)

    # Dark background banner on left
    banner = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5.8), Inches(7.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY_DARK
    banner.line.color.rgb = NAVY_DARK

    # Title & Metadata text
    t_box = s1.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = t_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "SCHOOL OF INTEGRATED ENGINEERING"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = BLUE_ACCENT
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "Academic Resourcing &\nFaculty Workload\nArchitecture"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "Carrying the Load Together: Building a Transparent, Defensible Resourcing Model for Our Departments"
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_after = Pt(28)

    p3 = tf.add_paragraph()
    p3.text = "Discussion & Collaborative Design Briefing\nUSAFA School of Integrated Engineering Leadership"
    p3.font.size = Pt(11)
    p3.font.color.rgb = RGBColor(148, 163, 184)

    # Hero Image on Right
    if os.path.exists(IMG_CARRYING):
        s1.shapes.add_picture(IMG_CARRYING, Inches(5.8), Inches(0), Inches(7.533), Inches(7.5))

    set_speaker_notes(s1, """GOOD MORNING, DEAN, DEPARTMENT HEADS, AND COLLEAGUES.

Welcome. We are gathered today to initiate an essential conversation about how we measure, balance, and advocate for our most vital resource across the School of Integrated Engineering: our faculty's time and talent.

I want to be clear about the purpose of this meeting from the first moment:
This is NOT an administrative audit.
This is NOT an effort to police faculty hours, tell any department they are "over-staffed," or enforce a rigid, one-size-fits-all formula.

Rather, this is an invitation to collaborate on an enduring advocacy tool. We want a shared, data-backed standard that helps you—as department leaders—defend your faculty lines, justify billets to the Dean and higher headquarters, balance internal burdens, and prevent burnout.

The image on your right represents the foundation of our conversation: across every department in our school, we are carrying a massive collective load. Today is about understanding all the pieces of that boulder and making sure our resourcing tools accurately reflect the work your faculty actually do.""")

    # =========================================================================
    # SLIDE 2: The Anatomy of Our Workload (The 5 Facets)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_header(s2, "The Anatomy of Departmental Workload", 
                 "Recognizing that Classroom Teaching is Only One Segment of Our Collective Mission")

    if os.path.exists(IMG_FACETS):
        s2.shapes.add_picture(IMG_FACETS, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5.1))

    facets = [
        ("1. Teaching & Curriculum", "Classroom lectures, lab courses, capstone senior design, course director duties, exam creation, and continuous grading.", NAVY_MID),
        ("2. Department Administration", "Department headship, executive officer duties, scheduling, personnel management, ABET accreditation, and hiring.", NAVY_DARK),
        ("3. Research & Scholarship", "Sponsored research execution, AFRL/NASA collaboration, faculty publications, conference leadership, and cadet research.", BLUE_ACCENT),
        ("4. Lab Operations & Safety", "Managing wind tunnels, machine shops, cyber ranges, high-hazard safety compliance, and technician oversight.", AMBER_ACCENT),
        ("5. Cadet Development & Institutional Service", "Official academic advising, ARC / Honor boards, cadet squadron mentorship, and Academy military support.", GREEN_ACCENT),
    ]

    card_top = 1.9
    card_height = 0.92
    card_gap = 0.12
    for i, (f_title, f_desc, f_color) in enumerate(facets):
        top_pos = card_top + i * (card_height + card_gap)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(top_pos), Inches(5.8), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        strip = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(top_pos), Inches(0.15), Inches(card_height))
        strip.fill.solid()
        strip.fill.fore_color.rgb = f_color
        strip.line.fill.background()

        tb = s2.shapes.add_textbox(Inches(7.0), Inches(top_pos + 0.08), Inches(5.35), Inches(card_height - 0.16))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p1 = tf.paragraphs[0]
        p1.text = f_title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = f_color
        p1.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = f_desc
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_MAIN

    set_speaker_notes(s2, """WHEN WE EXAMINE THE REALITY OF OUR FACULTY'S LIVES, THAT BOULDER IS NOT A MONOLITH.

It is comprised of five distinct, interlocking facets:

1. TEACHING & CURRICULUM: Preparing for 40-lesson courses, running laboratories, mentoring capstone teams, writing and grading exams.
2. DEPARTMENT ADMINISTRATION: Running a department, directing courses, managing personnel, handling ABET accreditation reviews, building class schedules.
3. RESEARCH & SCHOLARSHIP: Carrying sponsored grants, publishing, collaborating on defense initiatives, and driving the intellectual vitality of the school.
4. LAB OPERATIONS & SAFETY: Especially critical in the School of Integrated Engineering—maintaining wind tunnels, engine test cells, machining facilities, cyber ranges, and heavy equipment safely.
5. CADET DEVELOPMENT & SERVICE: Unique to our military academy mission—academic advising, serving on Academic Review (ARC) and Honor boards, serving as squadron liaisons, and supporting summer military training.

WHY DO WE HIGHLIGHT THIS UP FRONT?
Because any tool that ONLY looks at teaching and ignores the other four facets will produce a distorted, unfair, and demoralizing picture of your department.

Our ultimate vision is to account for the FULL boulder. Today, we invite your guidance on how to define and measure all five pieces.""")

    # =========================================================================
    # SLIDE 3: Different Roles, Different Proportions (The Tiered Model)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_header(s3, "Different Roles, Different Proportions", 
                 "Why Flat 'One-Size-Fits-All' Expectations Distort Department Reality")

    if os.path.exists(IMG_MULTITASK):
        s3.shapes.add_picture(IMG_MULTITASK, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5.1))

    tiers = [
        ("Department Leadership & Lab Directors", "Target: ~1 Section / Semester",
         "Department Heads, Deans, and major high-hazard Lab Directors carry immense administrative and operational responsibility. Expecting 3 sections is unrealistic and compromises safety and governance.", NAVY_DARK),
        ("Course Directors, Discipline Leads, & Execs", "Target: ~2 Sections / Semester",
         "Faculty leading multi-section core courses, managing ABET sub-disciplines, or serving as department executives require course release to coordinate instruction effectively.", NAVY_MID),
        ("Line Teaching Faculty", "Target: ~3 Sections / Semester",
         "Core instructional faculty whose primary operational allocation is classroom contact, cadet grading, and office hours.", BLUE_ACCENT),
        ("Endowed Chairs & Agency MOA Faculty", "Custom / Courtesy Allocation",
         "Visiting endowed chairs (part-time or allocated across departments by need) and MOA research partners (e.g., AFRL and NASA researchers like Dr. EH and Dr. HP teaching 1 section as a courtesy).", AMBER_ACCENT)
    ]

    top_pos = 1.9
    card_h = 1.15
    gap_h = 0.12
    for i, (t_title, t_badge, t_desc, t_col) in enumerate(tiers):
        cur_top = top_pos + i * (card_h + gap_h)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(cur_top), Inches(5.8), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        strip = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(cur_top), Inches(0.15), Inches(card_h))
        strip.fill.solid()
        strip.fill.fore_color.rgb = t_col
        strip.line.fill.background()

        tb = s3.shapes.add_textbox(Inches(7.0), Inches(cur_top + 0.08), Inches(5.35), Inches(card_h - 0.16))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = t_title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        p1.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = f"Expected Baseline: {t_badge}"
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.color.rgb = t_col
        p2.space_after = Pt(3)

        p3 = tf.add_paragraph()
        p3.text = t_desc
        p3.font.size = Pt(9)
        p3.font.color.rgb = TEXT_MAIN

    set_speaker_notes(s3, """JUST AS THE BOULDER HAS DIFFERENT PIECES, EACH INDIVIDUAL FACULTY MEMBER HAS A DIFFERENT CAPACITY PROFILE.

Think of the multi-armed faculty member shown here. One person is balancing a course lecture, a research proposal, a budget review, and a cadet honor hearing all at once.

AND HERE IS THE CORE PROBLEM WITH TRADITIONAL WORKLOAD REPORTS:
When an analytical model assumes that every single faculty member has a flat expectation of teaching 3 sections per semester:
- Departments with heavy lab oversight appear "under-loaded."
- Departments whose leaders teach 1 section appear to have "excess capacity."
- Departments that host specialized research fellows—like AFRL or NASA researchers who teach 1 courtesy section—are penalized because those instructors pull down the department's "average sections per instructor."

AS OUR DEAN HAS RIGHTLY POINTED OUT:
We need business rules that reflect reality:
- Deans, Department Heads, and major Lab Directors should target ~1 section.
- Course Directors, Discipline Leads, and Execs should target ~2 sections.
- Line instructional faculty target ~3 sections.
- Endowed Chairs and MOA adjuncts should be discounted and counted as bonus/fractional capacity.

Our goal is to build these tiered baselines directly into the tool, so when we evaluate whether a department is over- or under-capacity, we are measuring against REAL expectations, not a fictional standard.""")

    # =========================================================================
    # SLIDE 4: Phase 1: Focusing on the Teaching Core
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_header(s4, "Phase 1: Focusing on the Teaching Core", 
                 "Establishing Our First Objective Baseline Where Granular Institutional Data Exists")

    if os.path.exists(IMG_SPOTLIGHT):
        s4.shapes.add_picture(IMG_SPOTLIGHT, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5.1))

    reasons = [
        ("1. Objective & Granular Data Source",
         "The Registrar database records every course, section, cadet enrollment, and instructor assignment with high fidelity. Unlike research or committee hours, teaching enrollment data is clean, deterministic, and verifiable.", NAVY_MID),
        ("2. Direct Mission Impact to Cadets",
         "Classroom contact hours and student grading volume represent our primary direct interaction with the Cadet Wing. Establishing a sound baseline here ensures cadet academic quality is protected.", NAVY_DARK),
        ("3. Analytical Proof of Concept",
         "By perfecting our accounting engine on teaching first—solving multi-instructor splitting, duration weights, and cross-teaching—we establish a validated framework ready to ingest research, lab ops, and service data.", BLUE_ACCENT),
    ]

    top_pos = 1.9
    card_h = 1.55
    gap_h = 0.2
    for i, (r_title, r_desc, r_col) in enumerate(reasons):
        cur_top = top_pos + i * (card_h + gap_h)
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(cur_top), Inches(5.8), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        strip = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(cur_top), Inches(0.15), Inches(card_h))
        strip.fill.solid()
        strip.fill.fore_color.rgb = r_col
        strip.line.fill.background()

        tb = s4.shapes.add_textbox(Inches(7.0), Inches(cur_top + 0.12), Inches(5.35), Inches(card_h - 0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = r_title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        p1.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = r_desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MAIN

    set_speaker_notes(s4, """WHY DID WE START WITH TEACHING?

As shown in this spotlight illustration, we set the 'Teaching' block on our analysis table first—not because it is more important than research, administration, or labs—but because the data is readily available, highly structured, and verifiable across all departments.

1. OBJECTIVE DATA: We have complete semester-by-semester records of every section, cadet, and instructor.
2. DIRECT CADET TOUCHPOINT: Classroom teaching represents our core contact with future Air Force and Space Force officers.
3. OUR FOUNDATIONAL PROOF OF CONCEPT: If we can agree on how to count teaching fairly, we can apply that exact same transparent logic to research, lab management, and administrative burdens.

Let's look at how we built this initial teaching engine.""")

    # =========================================================================
    # SLIDE 5: Data Provenance & Accounting Rules
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_header(s5, "Data Provenance & The Accounting Model", 
                 "Transforming Raw Registrar Enrollments into Normalized, Fair Workload Metrics")

    col_w = 3.65
    col_gap = 0.35
    col_h = 5.0
    steps = [
        ("STEP 1: Granular Ingestion",
         "Direct Database Query",
         NAVY_DARK,
         [
             ("Query Source", "Registrar enrollment database for multi-semester tracking (Terms 2251, 2258, 2261, 2268)."),
             ("Fields Ingested", "Term, Class Number, Subject, Course Number, Course Title, Section Code, Cadet EMPLID, Instructor Name(s)."),
             ("Integrity Checks", "Regex parsing for multi-instructor strings, suffixes (Jr., III, Ph.D.), unassigned tokens (TBD, Staff)."),
         ]),
        ("STEP 2: Co-Teaching Attribution",
         "Fair Split Model (1/k)",
         NAVY_MID,
         [
             ("The Problem", "When two instructors co-teach a 24-cadet section, standard reports either double-count both (claiming 48 cadets taught) or credit only the primary instructor."),
             ("Our Solution", "Burden is split proportionally: each receives (Section Weight / k) and (Cadet Count / k)."),
             ("Audit Trail", "Unique individual cadet IDs and total seat counts are preserved for verification."),
         ]),
        ("STEP 3: Course Duration Weighting",
         "Semester-Equivalent Normalization",
         BLUE_ACCENT,
         [
             ("Full Semester (1.0 sec / 1.0 stu)", "Standard 40-lesson lecture/lab course (e.g. MECHENGR 330, MATH 141)."),
             ("Half Semester (0.5 sec / 0.5 stu)", "20-lesson courses (e.g. COMMSTRT 101). Two sections = 1 full-semester load."),
             ("Quarter Blocks (0.25 sec / 0.25 stu)", "10-lesson modular blocks (A, B, C, D sections). Four blocks = 1 full-semester load."),
             ("Independent Study 499 (0.0 sec / 1.0 stu)", "1-on-1 directed research. Zero classroom lecture prep, but full 1.0 contact credit per cadet."),
         ]),
    ]

    for i, (s_title, s_sub, s_col, items) in enumerate(steps):
        left_pos = 0.8 + i * (col_w + col_gap)
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.9), Inches(col_w), Inches(col_h))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        top_strip = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_pos), Inches(1.9), Inches(col_w), Inches(0.12))
        top_strip.fill.solid()
        top_strip.fill.fore_color.rgb = s_col
        top_strip.line.fill.background()

        tb = s5.shapes.add_textbox(Inches(left_pos + 0.2), Inches(2.15), Inches(col_w - 0.4), Inches(col_h - 0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = s_title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        p1.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = s_sub
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.color.rgb = s_col
        p2.space_after = Pt(14)

        for item_h, item_b in items:
            ph = tf.add_paragraph()
            ph.text = f"• {item_h}:"
            ph.font.size = Pt(10)
            ph.font.bold = True
            ph.font.color.rgb = NAVY_DARK
            ph.space_after = Pt(1)

            pb = tf.add_paragraph()
            pb.text = f"  {item_b}"
            pb.font.size = Pt(9.5)
            pb.font.color.rgb = TEXT_MAIN
            pb.space_after = Pt(8)

    set_speaker_notes(s5, """HERE IS THE EXACT DATA PIPELINE AND MATHEMATICAL ENGINE WE BUILT:

STEP 1: INGESTION
We pull raw enrollment data from the registrar. Each row represents a single cadet enrolled in a specific class number. Our engine cleans names, splits co-instructor strings, and removes unassigned/blank placeholder codes.

STEP 2: CO-TEACHING ATTRIBUTION (THE 1/k SPLIT)
A major distortion in previous workload reports was co-teaching. If Professor Smith and Major Jones co-teach a 24-student senior design section:
- In old spreadsheets, both got credit for 24 students (double-counting 48 student seats across the institution), OR one got 24 and the other got 0.
- In our engine, we split both the section prep burden and student contact load by 1/k. Each instructor receives 0.5 sections and 12 students.

STEP 3: DURATION WEIGHTING (SEMESTER-EQUIVALENT LOAD)
Not all sections demand the same preparation:
- Full-semester courses: 1.0 section prep and 1.0 student contact.
- Half-semester courses (like COMMSTRT 101): weighted at 0.5.
- Quarter-semester modular blocks: weighted at 0.25.
- And crucially, Independent Studies (499s): We award 0.0 section prep (no classroom lecture prep), but 1.0 full student contact load per cadet mentored. Mentoring 3 cadet researchers gives you 3.0 students of contact credit.

This ensures nobody gets credit for teaching a full lecture when they are mentoring an independent study, but they DO get full credit for their contact hours.""")

    # =========================================================================
    # SLIDE 6: Beyond the Averages: Section Distribution & Billet Reality
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_header(s6, "Beyond the Averages: The Sub-10 Cadet Dilemma", 
                 "Why Average Section Sizes Mask Elective Proliferation, Faculty Pain, and Billet Needs")

    chart_data = CategoryChartData()
    chart_data.categories = ['AERO', 'COMPSCI', 'ECE', 'MECH', 'DATASCI', 'CIVENGR', 'ASTRO', 'SYSENGR', 'ENGR', 'CYBER']
    chart_data.add_series('Sections with <= 10 Cadets', (10, 6, 17, 7, 3, 10, 1, 13, 3, 4))
    chart_data.add_series('Sections with > 10 Cadets', (61, 51, 27, 28, 32, 22, 27, 6, 8, 4))

    chart_frame = s6.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.8), Inches(1.9), Inches(6.5), Inches(5.1), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.size = Pt(9.5)
    chart.value_axis.has_major_gridlines = True

    top_pos = 1.9
    card_w = 4.9
    card_h = 1.55
    gap_h = 0.2
    insights = [
        ("The 'Average Class Size' Trap",
         "Department averages across our School range from ~13 to 25 cadets per section. But looking only at the mean conceals a major operational divide: high-enrollment core courses offset by heavily fragmented upper-level sections.",
         AMBER_ACCENT),
        ("Real Data: 74 Sections with <= 10 Cadets",
         "Across 340 engineering sections this semester, 74 (21.8%) have <= 10 cadets. In SYSENGR, 68.4% (13/19) have <= 10 cadets; CYBERSCI is 50.0% (4/8); ECE is 38.6% (17/44); CIVENGR is 31.2% (10/32); MECHENGR is 20.0% (7/35).",
         NAVY_MID),
        ("The Capstone & Elective Driver",
         "Granular course inspection reveals these small sections are NOT 499 independent studies—they are predominantly Senior Capstone Design sections (ECE 463, CIVENGR 451, SYSENGR 491) and specialized depth electives.",
         NAVY_DARK),
    ]

    for i, (in_title, in_desc, in_col) in enumerate(insights):
        cur_top = top_pos + i * (card_h + gap_h)
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(cur_top), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        strip = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(cur_top), Inches(0.15), Inches(card_h))
        strip.fill.solid()
        strip.fill.fore_color.rgb = in_col
        strip.line.fill.background()

        tb = s6.shapes.add_textbox(Inches(7.9), Inches(cur_top + 0.1), Inches(card_w - 0.45), Inches(card_h - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = in_title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        p1.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = in_desc
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_MAIN

    set_speaker_notes(s6, """OUR DEAN RAISED A CRITICAL QUESTION THAT GOES TO THE HEART OF RESOURCING:

He noted: 'ECE and CE look great on paper with ~13 to 16 cadets per section, but both departments claim to be hurting for faculty. What does the distribution look like, especially sections with 10 cadets or below?'

WE PULLED THIS SEMESTER'S ACTUAL REGISTRAR DATA (TERM 2268), AND THE RESULTS ARE EYE-OPENING:
Look at the chart on the left, covering all 340 engineering sections taught this term across our School:
- Across the entire school, 74 sections—nearly 22%—have 10 cadets or fewer.
- In SYSENGR: 13 out of 19 sections—over 68%—have 10 or fewer cadets!
- In CYBERSCI: 4 out of 8 sections (50%) have 10 or fewer cadets!
- In ECE: 17 out of 44 sections (nearly 39%) have 10 or fewer cadets!
- In CIVENGR: 10 out of 32 sections (31%) have 10 or fewer cadets!
- In MECHENGR: 7 out of 35 sections (20%) have 10 or fewer cadets!
- In AEROENGR: 10 out of 71 sections (14%) have 10 or fewer cadets!

AND CRUCIALLY: WHAT ARE THESE COURSES?
When we drilled down into the course numbers, they are NOT 1-on-1 independent studies. They are primarily:
1. Multi-section Senior Capstone projects—like ECE 463 (which has 9 sections with 3 to 7 cadets each), CIVENGR 451 (which has 7 small sections), and SYSENGR 491 (which has 13 small sections).
2. Upper-level depth electives (like ECE 434, 444, 447, CIVENGR 464, 485).

WHAT DOES THIS MEAN FOR DEPARTMENT HEADS?
1. THE AVERAGE IS MISLEADING: When you average a 28-cadet core class with two 5-cadet depth sections, your average is 13. It looks manageable on a high-level summary, but your faculty are burning 3 full course preparations.
2. CURRICULUM EFFICIENCY: It prompts an honest conversation: Can we consolidate capstone sections, rotate electives on alternate years, or pool cross-discipline design teams?
3. BILLET REALITY: It explains why faculty feel overwhelmed even when student ratios look balanced.

Our tool highlights these exact distributions so we can distinguish between real hiring need and elective consolidation opportunities.""")

    # =========================================================================
    # SLIDE 7: What the Baseline Shows: 2x2 Resourcing Quadrant Matrix
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_header(s7, "Department Resourcing Matrix (Initial Baseline)", 
                 "Balancing Course Preparation Burden (X-Axis) vs. Cadet Contact Burden (Y-Axis)")

    quad_left = 0.8
    quad_top = 1.9
    quad_w = 5.8
    quad_h = 5.1

    q_box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(quad_left), Inches(quad_top), Inches(quad_w), Inches(quad_h))
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = BG_LIGHT
    q_box.line.color.rgb = BORDER_COLOR

    q1 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(quad_left + 0.1), Inches(quad_top + 0.1), Inches(quad_w/2 - 0.15), Inches(quad_h/2 - 0.15))
    q1.fill.solid()
    q1.fill.fore_color.rgb = CARD_BG
    q1.line.color.rgb = BORDER_COLOR
    tf1 = q1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "QUADRANT 2: CORE LECTURES\nLow Preps / High Cadets\n(High grading, low prep friction)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    q2 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(quad_left + quad_w/2 + 0.05), Inches(quad_top + 0.1), Inches(quad_w/2 - 0.15), Inches(quad_h/2 - 0.15))
    q2.fill.solid()
    q2.fill.fore_color.rgb = AMBER_LIGHT
    q2.line.color.rgb = RGBColor(245, 158, 11)
    tf2 = q2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "QUADRANT 1: HIGH BURDEN\nHigh Preps / High Cadets\n(Acute Burnout Risk)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT

    q3 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(quad_left + 0.1), Inches(quad_top + quad_h/2 + 0.05), Inches(quad_w/2 - 0.15), Inches(quad_h/2 - 0.15))
    q3.fill.solid()
    q3.fill.fore_color.rgb = CARD_BG
    q3.line.color.rgb = BORDER_COLOR
    tf3 = q3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "QUADRANT 3: CAPACITY AVAILABLE\nLow Preps / Low Cadets\n(Or uncaptured lab/research load)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = TEXT_MUTED

    q4 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(quad_left + quad_w/2 + 0.05), Inches(quad_top + quad_h/2 + 0.05), Inches(quad_w/2 - 0.15), Inches(quad_h/2 - 0.15))
    q4.fill.solid()
    q4.fill.fore_color.rgb = CARD_BG
    q4.line.color.rgb = BORDER_COLOR
    tf4 = q4.text_frame
    tf4.word_wrap = True
    p = tf4.paragraphs[0]
    p.text = "QUADRANT 4: SPECIALIZED LABS\nHigh Preps / Low Cadets\n(Heavy prep, small class sizes)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NAVY_MID

    r_left = 7.0
    r_w = 5.5

    v_card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(r_left), Inches(1.9), Inches(r_w), Inches(2.7))
    v_card.fill.solid()
    v_card.fill.fore_color.rgb = CARD_BG
    v_card.line.color.rgb = BORDER_COLOR
    tb_v = s7.shapes.add_textbox(Inches(r_left + 0.2), Inches(2.05), Inches(r_w - 0.4), Inches(2.4))
    tf_v = tb_v.text_frame
    tf_v.word_wrap = True
    tf_v.margin_left = tf_v.margin_top = tf_v.margin_right = tf_v.margin_bottom = 0

    p1 = tf_v.paragraphs[0]
    p1.text = "The 2x2 Resourcing Matrix Concept"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_DARK
    p1.space_after = Pt(6)

    bullets = [
        "X-Axis (Horizontal): Weighted Sections / Instructor (Course Preparation Load).",
        "Y-Axis (Vertical): Cadets / Instructor (Direct Student Contact & Grading Load).",
        "Bubble Size: Total Cadets Enrolled in the Department.",
        "Color Gradient: Visual alerts moving from Green (balanced) to Amber to Red (high stress)."
    ]
    for b in bullets:
        pb = tf_v.add_paragraph()
        pb.text = f"• {b}"
        pb.font.size = Pt(9.5)
        pb.font.color.rgb = TEXT_MAIN
        pb.space_after = Pt(3)

    demo_box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(r_left), Inches(4.8), Inches(r_w), Inches(2.2))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = BLUE_LIGHT
    demo_box.line.color.rgb = BLUE_ACCENT
    tb_d = s7.shapes.add_textbox(Inches(r_left + 0.25), Inches(4.95), Inches(r_w - 0.5), Inches(1.9))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = 0

    pd1 = tf_d.paragraphs[0]
    pd1.text = "⚙️ LIVE DEMONSTRATION & SETTINGS TAB"
    pd1.font.size = Pt(12)
    pd1.font.bold = True
    pd1.font.color.rgb = NAVY_MID
    pd1.space_after = Pt(6)

    pd2 = tf_d.add_paragraph()
    pd2.text = "Our dashboard includes a real-time Administrator & Settings engine. In the live tool, you can:\n" \
               "1. Toggle Faculty Attribution (Home Department Rollup vs. Course-level).\n" \
               "2. Adjust Course Weights live (Independent Study, Half-Semesters, Quarters).\n" \
               "3. Instantly watch the 2x2 matrix and rankings recalculate in your browser.\n\n" \
               "You have the steering wheel—the model reflects YOUR agreed rules."
    pd2.font.size = Pt(9.5)
    pd2.font.color.rgb = NAVY_DARK

    set_speaker_notes(s7, """HERE IS HOW WE BRING THIS TOGETHER VISUALLY IN OUR EXECUTIVE DASHBOARD.

Rather than looking at a spreadsheet with 50 columns, we plot departments on a 2x2 Resourcing Matrix:
- HORIZONTAL (X-AXIS): PREPARATION LOAD. How many distinct, weighted sections is each faculty member preparing for?
- VERTICAL (Y-AXIS): CONTACT LOAD. How many cadets is each faculty member actually teaching, mentoring, and grading?

LOOK AT THE FOUR QUADRANTS:
- Top Right: High Preps + High Cadets. These departments are in acute danger of burnout.
- Top Left: High Cadets, but Low Preps. Typically core lecture courses with shared syllabi.
- Bottom Right: High Preps, but Low Cadets. This is where engineering specialized electives and heavy labs live—lots of faculty time spent preparing for small cadet cohorts.
- Bottom Left: Low Cadets, Low Preps. If a department is here, either they have surplus teaching capacity, OR more likely, THEIR LAB, RESEARCH, AND ADMINISTRATIVE LOAD IS NOT CURRENTLY BEING CAPTURED!

[PAUSE FOR A 2-MINUTE LIVE DEMO OF THE HTML DASHBOARD]
At this point in the meeting, I open the interactive dashboard file in a browser. I show them the 2x2 matrix, click on the Settings tab, change an Independent Study weight or toggle Primary Home Department, and show them how the entire dashboard recalculates live.

The key message: 'You have the steering wheel. We did not build a rigid black box; we built an engine that adapts to our school's consensus.'""")

    # =========================================================================
    # SLIDE 8: Accounting for Non-Teaching Burdens
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_header(s8, "The Rest of the Boulder: Accounting for Non-Teaching Loads", 
                 "Open Leadership Discussion: How Do We Fairly Quantify and Credit Non-Classroom Contributions?")

    card_w = 5.65
    card_h = 2.35
    row_top_1 = 1.9
    row_top_2 = 4.5
    col_left_1 = 0.8
    col_left_2 = 6.85

    d_cards = [
        ("1. Department Administration & Leadership",
         col_left_1, row_top_1, NAVY_DARK,
         [
             "Can we assume a flat administrative baseline across departments, or does burden scale with major size?",
             "How do we formally credit Course Directors who coordinate 10+ sections and junior instructors?",
             "Accounting for major governance roles: ABET accreditation leads, hiring committee chairs, schedulers."
         ]),
        ("2. Research, Scholarship, & MOA Collaborations",
         col_left_2, row_top_1, BLUE_ACCENT,
         [
             "How do we capture research load without turning this into a grant dollar competition?",
             "Should we measure faculty FTE time allocation (e.g. 20% research track) rather than outputs?",
             "How do we credit AFRL and NASA agency researchers (Dr. EH, Dr. HP) as fractional/courtesy teaching?"
         ]),
        ("3. Lab Operations, Facilities, & Safety",
         col_left_1, row_top_2, AMBER_ACCENT,
         [
             "Engineering labs have massive physical footprints: wind tunnels, engine test cells, cyber ranges.",
             "What metrics reflect lab load? Annual student user-hours? Equipment maintenance value? Safety risk?",
             "Ensuring Lab Directors and technicians receive proper workload credit for maintaining operational readiness."
         ]),
        ("4. Cadet Advising & Institutional Service",
         col_left_2, row_top_2, GREEN_ACCENT,
         [
             "Cadet advising numbers are tracked in the database: where does advising load fit into the boulder?",
             "Service academy military duties: Squadron mentorship, AOC liaison, summer training (Ops AF, Field Engr).",
             "Academic governance: Academic Review Committee (ARC) boards, Honor boards, faculty senate."
         ]),
    ]

    for (c_title, c_left, c_top, c_col, c_bullets) in d_cards:
        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c_left), Inches(c_top), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        strip = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(c_left), Inches(c_top), Inches(card_w), Inches(0.12))
        strip.fill.solid()
        strip.fill.fore_color.rgb = c_col
        strip.line.fill.background()

        tb = s8.shapes.add_textbox(Inches(c_left + 0.2), Inches(c_top + 0.2), Inches(card_w - 0.4), Inches(card_h - 0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = c_title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        p1.space_after = Pt(6)

        for b in c_bullets:
            pb = tf.add_paragraph()
            pb.text = f"• {b}"
            pb.font.size = Pt(9.5)
            pb.font.color.rgb = TEXT_MAIN
            pb.space_after = Pt(4)

    set_speaker_notes(s8, """NOW WE RETURN TO THE REST OF THE BOULDER—AND THIS IS WHERE WE NEED YOUR INPUT.

We want to open the floor to discuss how we capture the non-teaching areas:

1. ADMINISTRATION: Can we assume all departments have the same administrative burden? A department with 200 majors and ABET accreditation probably carries more administrative friction than a small minor. How do we account for course directors?
2. RESEARCH: We do NOT want a punitive grant-dollar competition. A better approach is measuring CAPACITY ALLOCATION: Is this faculty member budgeted for 20% research, 40% research, or on a teaching track? And how do we credit our agency MOA researchers?
3. LAB OPERATIONS: In the School of Integrated Engineering, our labs are massive capital assets. How should we measure lab burden—by student contact hours in labs? By equipment complexity? By safety oversight?
4. CADET ADVISING & ACADEMY SERVICE: We have registrar data on who advises each cadet. Mentoring 30 advisees through academic probation and ARC boards is real work. Where does that count?

Let's hear your thoughts on which metrics make sense and which ones we should avoid.""")

    # =========================================================================
    # SLIDE 9: Our Collaborative Request: Department Starter Rosters
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_header(s9, "Our Collaborative Request: Department Starter Rosters", 
                 "Eliminating Blind Spots by Validating Faculty Roles, Billet Status, and Effort Allocations")

    left_card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(4.5), Inches(5.1))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG
    left_card.line.color.rgb = BORDER_COLOR
    tb_l = s9.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(4.1), Inches(4.7))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0

    p = tf_l.paragraphs[0]
    p.text = "The Department Starter Kit"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_DARK
    p.space_after = Pt(8)

    p_body = [
        ("No Blank-Sheet Homework", "We are NOT asking you to spend hours typing spreadsheets from scratch. We will provide each department head with a pre-populated roster."),
        ("Pre-Populated Data", "Contains your faculty names, verified teaching loads from the registrar, and cadet advisee counts."),
        ("What We Need From You", "A quick validation of four specific attributes for each faculty member on your roster."),
        ("Turnaround Goal", "Review, annotate exceptions, and return within 2 weeks so we can integrate your ground truth into the school model.")
    ]
    for h, b in p_body:
        ph = tf_l.add_paragraph()
        ph.text = f"• {h}:"
        ph.font.size = Pt(10)
        ph.font.bold = True
        ph.font.color.rgb = NAVY_MID
        pb = tf_l.add_paragraph()
        pb.text = f"  {b}"
        pb.font.size = Pt(9.5)
        pb.font.color.rgb = TEXT_MAIN
        pb.space_after = Pt(6)

    tbl_shape = s9.shapes.add_table(6, 4, Inches(5.6), Inches(1.9), Inches(6.9), Inches(4.0))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(1.8)

    headers = ["Faculty Member", "Billet Status", "Expected Tier", "Primary Non-Teaching Allocation"]
    for col_idx, h_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    sample_data = [
        ("Col J. Miller (DH)", "Filled (Military)", "Dept Head (1 sec)", "50% Dept Administration"),
        ("Dr. S. Williams (Prof)", "Filled (Civilian)", "Line Faculty (3 secs)", "15% Sponsored Research"),
        ("Lt Col D. Brown", "Double-Billeted", "Course Dir (2 secs)", "30% Aero Lab Direction"),
        ("Dr. E. Hansen (AFRL)", "MOA Adjunct", "Courtesy (1 sec)", "80% AFRL Hypersonics"),
        ("Vacant Billet #14", "Unfilled (Civilian Line)", "Line Faculty (3 secs)", "Hiring in Progress (Strain)")
    ]

    for row_idx, row_vals in enumerate(sample_data, start=1):
        for col_idx, val in enumerate(row_vals):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_LIGHT if row_idx % 2 == 0 else CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_MAIN

    b_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), Inches(6.1), Inches(6.9), Inches(0.9))
    b_box.fill.solid()
    b_box.fill.fore_color.rgb = AMBER_LIGHT
    b_box.line.color.rgb = AMBER_ACCENT
    tb_b = b_box.text_frame
    tb_b.word_wrap = True
    pb = tb_b.paragraphs[0]
    pb.text = "⭐ CAPTURING BILLET REALITY: Tracking vacant and double-billeted lines allows the Dean to understand whether a department's stress is temporary (PCS/hiring lag) or structural, supporting smart billet reallocation across the School."
    pb.font.size = Pt(9.5)
    pb.font.bold = True
    pb.font.color.rgb = RGBColor(146, 64, 14)

    set_speaker_notes(s9, """HERE IS OUR SPECIFIC REQUEST TO EACH DEPARTMENT HEAD:

We are NOT asking you to go write a 10-page report or do blank-sheet math. We will send each of you a 'Starter Roster' pre-populated with:
1. The list of instructors teaching in your discipline from registrar data.
2. Their section counts and student contact counts.
3. Their official advisee counts.

ALL WE NEED YOU TO DO IS VALIDATE FOUR COLUMNS:
1. BILLET STATUS: Is this billet filled, vacant, double-billeted, or a temporary PCS turnover? This is essential for the Dean to know where the shoe is really pinching.
2. EXPECTED TIER: Is this person a Department Head (1 sec), Course Director (2 secs), Line Faculty (3 secs), or an Endowed Chair/MOA?
3. PRIMARY HOME AFFILIATION: Confirm whether interdisciplinary instructors belong to you or another department.
4. ESTIMATED EFFORT SPLIT: A high-level estimate of major non-teaching allocations (e.g. 30% lab direction, 40% admin, 20% research).

This low-friction starter kit gives us the verified ground truth we need without burying you in administrative overhead.""")

    # =========================================================================
    # SLIDE 10: The Road Ahead: An Enduring Resourcing Platform
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    apply_header(s10, "The Road Ahead: An Enduring Resourcing Platform", 
                 "Transitioning from Subject Codes to Department-Centric Resourcing & Advocacy")

    m_w = 3.65
    m_gap = 0.35
    m_h = 4.2
    milestones = [
        ("MILESTONE 1: INGESTION",
         "Ingest Department Rosters",
         NAVY_DARK,
         [
             ("Timeline", "Next 2–3 Weeks"),
             ("Actions", "Distribute starter rosters to all engineering department heads."),
             ("Deliverables", "Clean, deduplicated faculty mapping by official Home Department (DFAN, DFAS, DFCE, DFCS, DFEC, DFEM)."),
             ("Billet Tracking", "Capture vacant, filled, and double-billeted lines across the entire School.")
         ]),
        ("MILESTONE 2: CALIBRATION",
         "Implement Tiered Expectations",
         NAVY_MID,
         [
             ("Timeline", "Mid-Semester Review"),
             ("Actions", "Apply Dean's tiered baseline rules (1, 2, 3 sections) to recalculate expected vs. actual teaching loads."),
             ("Non-Teaching Offsets", "Incorporate accredited lab management, course directorships, and research MOAs."),
             ("Section Health", "Flag elective proliferation and sub-10 cadet section consolidation opportunities.")
         ]),
        ("MILESTONE 3: ADVOCACY",
         "School Resourcing Platform",
         BLUE_ACCENT,
         [
             ("Timeline", "End of Term / Annual Planning"),
             ("Actions", "Deliver the comprehensive, privacy-secured executive dashboard to School leadership."),
             ("Outcomes", "Equitable faculty workload balancing across engineering disciplines."),
             ("Advocacy Power", "Data-backed evidence to defend faculty lines, hire tenure/civilian lines, and optimize cadet contact.")
         ]),
    ]

    for i, (m_tag, m_title, m_col, m_items) in enumerate(milestones):
        left_pos = 0.8 + i * (m_w + m_gap)
        card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.9), Inches(m_w), Inches(m_h))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        top_strip = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_pos), Inches(1.9), Inches(m_w), Inches(0.12))
        top_strip.fill.solid()
        top_strip.fill.fore_color.rgb = m_col
        top_strip.line.fill.background()

        tb = s10.shapes.add_textbox(Inches(left_pos + 0.2), Inches(2.15), Inches(m_w - 0.4), Inches(m_h - 0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = m_tag
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = m_col
        p1.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = m_title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = NAVY_DARK
        p2.space_after = Pt(14)

        for item_h, item_b in m_items:
            ph = tf.add_paragraph()
            ph.text = f"• {item_h}:"
            ph.font.size = Pt(10)
            ph.font.bold = True
            ph.font.color.rgb = NAVY_DARK
            ph.space_after = Pt(1)

            pb = tf.add_paragraph()
            pb.text = f"  {item_b}"
            pb.font.size = Pt(9.5)
            pb.font.color.rgb = TEXT_MAIN
            pb.space_after = Pt(8)

    bot_card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.85))
    bot_card.fill.solid()
    bot_card.fill.fore_color.rgb = NAVY_DARK
    bot_card.line.color.rgb = NAVY_DARK
    tb_bot = bot_card.text_frame
    tb_bot.word_wrap = True
    p_bot = tb_bot.paragraphs[0]
    p_bot.text = "OUR SHARED GOAL: An enduring, transparent platform that moves beyond anecdotal arguments—giving the School of Integrated Engineering a unified, defensible voice to secure the resources, billets, and faculty lines our mission demands."
    p_bot.font.size = Pt(10.5)
    p_bot.font.bold = True
    p_bot.font.color.rgb = RGBColor(255, 255, 255)
    p_bot.alignment = PP_ALIGN.CENTER

    set_speaker_notes(s10, """IN CLOSING, HERE IS OUR ROADMAP:

1. MILESTONE 1 (NEXT 2 WEEKS): We will send you your pre-populated Starter Rosters. You validate the home department assignments, flag your vacant/double-billeted lines, and assign the tiered expectations.
2. MILESTONE 2 (MID-SEMESTER): We ingest that data and re-run our models. Now, when we look at the School of Integrated Engineering, we aren't looking at generic subject codes; we are looking at real departments (DFAN, DFAS, DFCE, DFCS, DFEC, DFEM) evaluated against their true, tiered expected capacities.
3. MILESTONE 3 (ANNUAL RESOURCING PLATFORM): We deliver a finalized, executive dashboard. When our Dean goes to the Superintendent or higher headquarters to argue for civilian lines, military billets, or lab funding, he will have an irrefutable, transparent, school-wide dataset to back up every request.

Thank you for your time, your leadership, and your commitment to carrying this load together.

I look forward to your questions and discussion.""")

    prs.save(OUTPUT_PPTX)
    print(f"Successfully generated PowerPoint presentation: {os.path.abspath(OUTPUT_PPTX)}")


if __name__ == "__main__":
    build_presentation()
